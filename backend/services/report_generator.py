"""
Report Generator — Creates PDF and PPTX reports from analysis results.
"""

import os
from datetime import datetime
from config import settings
from utils.helpers import safe_json_loads


async def generate_report(workflow, report_type: str = "pdf") -> dict:
    """Generate a report from workflow results."""
    insight_result = safe_json_loads(workflow.insight_result, {})
    audit_result = safe_json_loads(workflow.audit_result, {})
    viz_result = safe_json_loads(workflow.visualization_result, {})

    executive_summary = insight_result.get("executive_summary", "Analysis complete.")
    insights = insight_result.get("insights", [])
    recommendations = insight_result.get("recommendations", [])

    if report_type == "pdf":
        file_path = await _generate_pdf(workflow, executive_summary, insights, recommendations)
    elif report_type == "pptx":
        file_path = await _generate_pptx(workflow, executive_summary, insights, recommendations)
    else:
        file_path = None

    return {
        "file_path": file_path,
        "content": {"executive_summary": executive_summary},
        "executive_summary": executive_summary,
        "insights": insights,
        "recommendations": recommendations,
    }


async def _generate_pdf(workflow, summary, insights, recommendations) -> str:
    """Generate a PDF report using reportlab."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch

        filename = f"report_{workflow.id[:8]}_{datetime.now().strftime('%Y%m%d')}.pdf"
        filepath = os.path.join(settings.REPORTS_DIR, filename)

        doc = SimpleDocTemplate(filepath, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # Title
        story.append(Paragraph("Data Analysis Report", styles["Title"]))
        story.append(Spacer(1, 0.3 * inch))
        story.append(Paragraph(f"Query: {workflow.prompt}", styles["Heading2"]))
        story.append(Spacer(1, 0.2 * inch))

        # Executive Summary
        story.append(Paragraph("Executive Summary", styles["Heading2"]))
        for line in summary.split("\n"):
            clean = line.replace("#", "").replace("*", "").strip()
            if clean:
                story.append(Paragraph(clean, styles["BodyText"]))
        story.append(Spacer(1, 0.2 * inch))

        # Key Insights
        story.append(Paragraph("Key Insights", styles["Heading2"]))
        for i, insight in enumerate(insights[:10], 1):
            text = f"{i}. [{insight.get('category', '')}] {insight.get('title', '')} — {insight.get('description', '')}"
            story.append(Paragraph(text, styles["BodyText"]))
        story.append(Spacer(1, 0.2 * inch))

        # Recommendations
        story.append(Paragraph("Recommendations", styles["Heading2"]))
        for i, rec in enumerate(recommendations[:10], 1):
            text = f"{i}. [{rec.get('priority', '').upper()}] {rec.get('action', '')} — {rec.get('rationale', '')}"
            story.append(Paragraph(text, styles["BodyText"]))

        doc.build(story)
        return filepath
    except ImportError:
        return None


async def _generate_pptx(workflow, summary, insights, recommendations) -> str:
    """Generate a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        filename = f"report_{workflow.id[:8]}_{datetime.now().strftime('%Y%m%d')}.pptx"
        filepath = os.path.join(settings.REPORTS_DIR, filename)

        prs = Presentation()
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Data Analysis Report"
        slide.placeholders[1].text = workflow.prompt

        # Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        body = slide.placeholders[1]
        clean_summary = summary.replace("#", "").replace("*", "")
        body.text = clean_summary[:800]

        # Insights slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Insights"
        body = slide.placeholders[1]
        text = "\n".join([f"• {i.get('title', '')}" for i in insights[:8]])
        body.text = text

        # Recommendations slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Recommendations"
        body = slide.placeholders[1]
        text = "\n".join([f"• [{r.get('priority', '').upper()}] {r.get('action', '')}" for r in recommendations[:8]])
        body.text = text

        prs.save(filepath)
        return filepath
    except ImportError:
        return None
